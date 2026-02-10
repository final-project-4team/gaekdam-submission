from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from PIL import Image, ImageOps
import pytesseract

def _ocr_image_pil(img: Image.Image, lang="kor+eng") -> str:
    # 간단 전처리 : 그레이스케일 + 리사이즈 + 이진화
    img = ImageOps.grayscale(img)
    w, h = img.size
    if max(w, h) < 1000:
        img = img.resize((int(w*2), int (h*2)), Image.LANCZOS)
    img = img.point(lambda p: 0 if p < 128 else 255)
    return pytesseract.image_to_string(img, lang=lang)

def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides):
        texts = []
        # 텍스트 모음
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
        
        # 이미지 -> OCR
        ocr_texts = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, "image"):
                blob = shape.image.blob
                try:
                    img = Image.open(BytesIO(blob)).convert("RGB")
                    txt = _ocr_image_pil(img)
                    if txt and txt.strip():
                        ocr_texts.append("[IMAGE_OCR]\n" + txt.strip())
                except Exception:
                    pass

        # 발표자 노트
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text or ""
        slide_text = "\n".join([t for t in texts if t])
        if notes and notes.strip():
            slide_text += "\n\n[NOTES]\n" + notes.strip()
        if ocr_texts:
            slide_text += "\n\n" + "\n\n".join(ocr_texts)
        if slide_text.strip():
            parts.append(f"[SLIDE {i+1}]\n{slide_text.strip()}")
    return "\n\n".join(parts)